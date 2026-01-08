#!/usr/bin/env python3
# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx",
# ]
# ///
"""
Generate alignment-hive lightning talk slides.

Run with: uv run generate_slides.py
Output: alignment-hive-talk.pptx (upload to Google Drive, open as Google Slides)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Colors (dark theme)
BG_COLOR = RGBColor(0x1a, 0x1a, 0x2e)  # Dark blue-black
TEXT_COLOR = RGBColor(0xea, 0xea, 0xea)  # Off-white
ACCENT_COLOR = RGBColor(0x6c, 0xb4, 0xee)  # Light blue
CODE_BG = RGBColor(0x2d, 0x2d, 0x44)  # Slightly lighter for code
MUTED_COLOR = RGBColor(0xaa, 0xaa, 0xaa)  # Gray for descriptions


def set_slide_background(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle=None):
    """Add a title slide with centered text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_background(slide, BG_COLOR)

    # Title
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = ACCENT_COLOR
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(20)

    return slide


def add_content_slide(prs, title, bullets, code=None, code_top_inches=4.2):
    """Add a content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_background(slide, BG_COLOR)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    # Bullets
    bullet_top = Inches(1.1)
    bullet_box = slide.shapes.add_textbox(
        Inches(0.5), bullet_top, Inches(9), Inches(3.8)
    )
    tf = bullet_box.text_frame
    tf.word_wrap = True

    first = True
    for bullet in bullets:
        # Skip empty bullets
        if not bullet.strip():
            continue
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.text = bullet
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(8)
        p.level = 0

    # Code block if provided
    if code:
        code_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(code_top_inches), Inches(8.5), Inches(0.6)
        )
        tf = code_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = code
        p.font.size = Pt(18)
        p.font.name = "Courier New"
        p.font.color.rgb = ACCENT_COLOR
        p.alignment = PP_ALIGN.LEFT

        # Add background to code (via shape fill)
        code_box.fill.solid()
        code_box.fill.fore_color.rgb = CODE_BG

    return slide


def add_cta_slide(prs, title, items):
    """Add a call-to-action slide with numbered items."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_COLOR)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    # Items
    item_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(4.2))
    tf = item_box.text_frame
    tf.word_wrap = True

    for i, (label, desc) in enumerate(items):
        if i > 0:
            p = tf.add_paragraph()
            p.space_before = Pt(14)
        else:
            p = tf.paragraphs[0]

        p.text = f"{i+1}. {label}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = TEXT_COLOR

        p = tf.add_paragraph()
        p.text = f"    {desc}"
        p.font.size = Pt(18)
        p.font.color.rgb = MUTED_COLOR

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 aspect ratio

    # Slide 1: Hook
    add_content_slide(
        prs,
        "Who are AI tools being built for?",
        [
            "AI tool builders optimize for:",
            "  - Their internal use case (dogfooding)",
            "  - Their largest revenue sources (corporations)",
            "",
            "By default, independent researchers get left behind.",
        ],
    )

    # Slide 2: The Mechanism
    add_content_slide(
        prs,
        "Why this happens",
        [
            "Big orgs build up: skills, plugins, instructions,",
            "accumulated knowledge.",
            "",
            "This gets distributed automatically to their AI agents.",
            "It's continuous learning in practice.",
            "",
            "Independents don't have this. The gap will widen.",
        ],
    )

    # Slide 3: Why This Matters
    add_content_slide(
        prs,
        "Soft takeoff has started. Now is the time.",
        [
            "AI is already accelerating research.",
            "Alignment work needs this acceleration too.",
            "",
            "The field is:",
            "  - Distributed across small teams & independents",
            "  - Working on urgent problems",
            "  - Not big enough to build custom infra",
        ],
    )

    # Slide 4: The Solution
    add_content_slide(
        prs,
        "alignment-hive",
        [
            "Shared infrastructure for alignment researchers",
            "",
            "  - Plugin marketplace for Claude Code",
            "  - Curated skills for common research workflows",
            "  - hive-mind: sharing session learnings (in dev)",
            "",
            "github.com/Crazytieguy/alignment-hive",
        ],
    )

    # Slide 5: What's Available Now
    add_content_slide(
        prs,
        "Start today",
        [
            "MATS plugin includes:",
            "  - project-setup: architecture decisions for new projects",
            "  - fellow-handbook: MATS policies, compute, housing",
            "",
            "One command to install:",
        ],
        code="/plugin install mats@alignment-hive",
        code_top_inches=2.9,
    )

    # Slide 6: Call to Action
    add_cta_slide(
        prs,
        "Let's build this together",
        [
            ("Try it", "/plugin install mats@alignment-hive"),
            ("Give feedback", "What's friction? What do you need?"),
            ("Pair with me", "Let me watch the tools in action"),
            ("Contribute", "Talk to me if interested"),
        ],
    )

    output_path = "alignment-hive-talk.pptx"
    prs.save(output_path)
    print(f"Created: {output_path}")
    print("Upload to Google Drive and open with Google Slides")


if __name__ == "__main__":
    main()
